"""
PPTX to Content Converter
=========================
Converts .pptx files to .pdf or .png/.jpg using the best available method.
"""

import sys, os, subprocess, shutil, logging
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

# Configure logging
logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')
logger = logging.getLogger(__name__)

class PPTXConverter:
    def __init__(self, libreoffice_path=None):
        self.soffice = libreoffice_path or self.find_libreoffice()
        self.has_powerpoint = self.check_powerpoint()

    def find_libreoffice(self):
        """Find LibreOffice installation."""
        candidates = [
            shutil.which("soffice"),
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            os.path.expanduser(r"~\AppData\Local\Programs\LibreOffice\program\soffice.exe"),
        ]
        for c in candidates:
            if c and Path(c).exists():
                return c
        return None 


    def check_powerpoint(self):
        """Check if PowerPoint is available via COM."""
        if os.name != 'nt': return False
        try:
            import comtypes.client
            # Just try to see if we can create the object
            pp = comtypes.client.CreateObject("PowerPoint.Application")
            pp.Quit()
            return True
        except:
            return False

    def convert_to_pdf(self, pptx_path, pdf_path):
        """Convert to PDF using the best available method."""
        pptx_path = Path(pptx_path).resolve()
        pdf_path = Path(pdf_path).resolve()
        
        if self.has_powerpoint:
            return self._convert_pdf_powerpoint(pptx_path, pdf_path)
        elif self.soffice:
            return self._convert_pdf_libreoffice(pptx_path, pdf_path)
        else:
            logger.warning("No PowerPoint or LibreOffice found. Using fallback extractor.")
            return self._convert_pdf_fallback(pptx_path, pdf_path)

    def _convert_pdf_powerpoint(self, pptx_path, pdf_path):
        import comtypes.client
        pp = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            # ppSaveAsPDF = 32
            deck = pp.Presentations.Open(str(pptx_path), WithWindow=False)
            deck.SaveAs(str(pdf_path), 32)
            deck.Close()
            return True
        except Exception as e:
            logger.error(f"PowerPoint conversion failed: {e}")
            return False
        finally:
            pp.Quit()

    def _convert_pdf_libreoffice(self, pptx_path, pdf_path):
        out_dir = str(pdf_path.parent)
        try:
            subprocess.run(
                [self.soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, str(pptx_path)],
                capture_output=True, text=True, timeout=120, check=True
            )
            # Rename if necessary (LibreOffice defaults to input_name.pdf)
            lo_output = pdf_path.parent / (pptx_path.stem + ".pdf")
            if lo_output.exists() and lo_output != pdf_path:
                if pdf_path.exists(): os.remove(pdf_path)
                lo_output.rename(pdf_path)
            return pdf_path.exists()
        except Exception as e:
            logger.error(f"LibreOffice conversion failed: {e}")
            return False

    def _convert_pdf_fallback(self, pptx_path, pdf_path):
        """Extract text and images into a simple PDF."""
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import landscape, A4
            from reportlab.lib.units import inch
        except ImportError:
            logger.error("Install dependencies for fallback: pip install python-pptx reportlab")
            return False

        try:
            prs = Presentation(pptx_path)
            c = canvas.Canvas(str(pdf_path), pagesize=landscape(A4))
            width, height = landscape(A4)

            for i, slide in enumerate(prs.slides, 1):
                c.setFont("Helvetica-Bold", 16)
                c.drawString(0.5*inch, height - 0.5*inch, f"Slide {i}")
                
                y = height - 1*inch
                c.setFont("Helvetica", 12)
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                # Simple wrapping
                                words = text.split()
                                line = ""
                                for word in words:
                                    if c.stringWidth(line + " " + word) < (width - 1*inch):
                                        line += " " + word
                                    else:
                                        c.drawString(0.7*inch, y, line.strip())
                                        y -= 15
                                        line = word
                                c.drawString(0.7*inch, y, line.strip())
                                y -= 20
                    
                    if y < 1*inch: break # Simple page overflow protection
                
                c.showPage()
            c.save()
            return True
        except Exception as e:
            logger.error(f"Fallback conversion failed: {e}")
            return False

    def convert_to_images(self, pptx_path, output_dir, format="png"):
        """Convert slides to individual images."""
        pptx_path = Path(pptx_path).resolve()
        output_dir = Path(output_dir).resolve()
        output_dir.mkdir(parents=True, exist_ok=True)

        # PowerPoint COM is best for this too if on Windows
        if self.has_powerpoint:
            return self._convert_images_powerpoint(pptx_path, output_dir, format)
        else:
            logger.warning("No PowerPoint found for image export. Using LibreOffice/Fallback.")
            # For now, we fallback to converting to PDF then PDF to Images if possible, 
            # or use python-pptx to get text/images (very limited)
            return self._convert_images_fallback(pptx_path, output_dir, format)

    def _convert_images_powerpoint(self, pptx_path, output_dir, format):
        import comtypes.client
        pp = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            deck = pp.Presentations.Open(str(pptx_path), WithWindow=False)
            # Export takes folder path and format
            # 17 = JPG, 18 = PNG
            pp_format = 18 if format.lower() == "png" else 17
            deck.Export(str(output_dir), format.upper()) 
            deck.Close()
            return True
        except Exception as e:
            logger.error(f"PowerPoint image export failed: {e}")
            return False
        finally:
            pp.Quit()

    def _convert_images_fallback(self, pptx_path, output_dir, format):
        """PDF to Images fallback (requires pdf2image and poppler)."""
        pdf_temp = output_dir / "temp_slides.pdf"
        if self.convert_to_pdf(pptx_path, pdf_temp):
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(pdf_temp))
                for i, img in enumerate(images):
                    img.save(output_dir / f"Slide{i+1}.{format}", format.upper())
                os.remove(pdf_temp)
                return True
            except ImportError:
                logger.error("Install pdf2image and Poppler for image conversion fallback.")
                return False
        return False

def main():
    import argparse
    parser = argparse.ArgumentParser(description="PPTX to PDF/Image Converter")
    parser.add_argument("input", help="Input .pptx or folder")
    parser.add_argument("-o", "--output", help="Output file or folder")
    parser.add_argument("-f", "--format", choices=["pdf", "png", "jpg"], default="pdf", help="Output format")
    args = parser.parse_args()

    converter = PPTXConverter()
    target = Path(args.input)

    if target.is_dir():
        files = list(target.glob("*.pptx"))
        logger.info(f"Found {len(files)} files in {target}")
        with ThreadPoolExecutor(max_workers=4) as executor:
            for f in files:
                out = Path(args.output or f.parent) / (f.stem + (f".{args.format}" if args.format == "pdf" else ""))
                if args.format == "pdf":
                    executor.submit(converter.convert_to_pdf, f, out)
                else:
                    executor.submit(converter.convert_to_images, f, out, args.format)
    else:
        out = Path(args.output or target.with_suffix(f".{args.format}" if args.format == "pdf" else ""))
        if args.format == "pdf":
            converter.convert_to_pdf(target, out)
        else:
            converter.convert_to_images(target, out, args.format)

if __name__ == "__main__":
    main()
